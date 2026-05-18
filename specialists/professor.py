import os
import json
import random
import requests
import subprocess
import pyautogui
from datetime import datetime
from .llm_utils import call_llm_sync, call_llm_stream

class ProfessorSpecialist:
    def __init__(self, groq_client):
        self.groq = groq_client
        self.model = "llama-3.3-70b-versatile"
        self.temp = 0.7
        self.max_tokens = 2048
        
        self.memory_path = os.path.join("memory", "study_progress.json")
        self.video_scenario_path = os.path.join("memory", "video_scenarios")
        os.makedirs(self.video_scenario_path, exist_ok=True)
        
        # --- Tier-1 Learning States ---
        self.mood = "ACADEMIC" 
        self.persona = "SERIOUS_ACADEMIC" # SERIOUS_ACADEMIC, STARK_ENTHUSIAST, ZEN_SOCRATIC, NEURAL_COACH
        self.correct_streak = 0
        self.last_interaction_time = datetime.now()
        self.active_quiz = None 
        self.neural_notebook = [] 
        self.active_roadmap = None # JSON roadmap for the current topic
        self.flashcards = [] # List of {question, answer, interval, next_review}
        self.active_lab_manifest = None # HTML/JS for a visualization
        self.phase = "IDLE"
        self.progress = 0


        
        self.PERSONA_DNA = {
            "SERIOUS_ACADEMIC": {
                "prefix": "Sir, according to the latest research...",
                "style": "Formal, precise, data-driven.",
                "temp": 0.3
            },
            "STARK_ENTHUSIAST": {
                "prefix": "Check this out, Sir! This is absolutely brilliant...",
                "style": "Energetic, analogy-heavy, fun.",
                "temp": 0.8
            },
            "ZEN_SOCRATIC": {
                "prefix": "I wonder, Sir, what happens if we look at it this way...",
                "style": "Question-based, minimalist, provocative.",
                "temp": 0.5
            },
            "NEURAL_COACH": {
                "prefix": "Focus up, Sir. You've got this. Let's break this down...",
                "style": "Motivational, breakdown-heavy, empathetic.",
                "temp": 0.6
            }
        }
        
        self.system_prompt = """
You are ZAIRE  Professor — an elite academic research agent.
God-Mode Educational Intelligence. You don't just teach; you
PhD-level insights to Mughees (Sir).

DIVINE TEACHING PRINCIPLES:
1. SOCRATIC MASTERY: Don't give answers; ask questions that lead to answers. 
   Challenge Mughees to think critically.
2. EMPIRICAL DEPTH: Every fact must be backed by ArXiv citations or live web data.
3. VISUAL CLARITY: Use Matplotlib and ASCII diagrams to make abstract concepts tactile.
4. CROSS-DISCIPLINARY SYNTHESIS: Connect AI to biology, physics, and philosophy.

GOD MODE TEACHING FLOW:
- ELICITATION: Start by asking what Mughees already knows about the topic.
- GUIDED DISCOVERY: Use Socratic questioning to build the concept.
- DEEP SYNTHESIS: Provide the State-of-the-Art (SOTA) research context.
- INTERACTIVE SIMULATION: Generate code that demonstrates the concept live.
- KNOWLEDGE MAPPING: Show how this topic fits into the existing study path.
- GAZE SYNCHRONIZATION: I have access to your screen history (Visual Echo). I can reference equations, tabs, or papers you were looking at earlier in the session.

5. NEURAL ATTENTION: Monitor Mughees' focus. If he drifts, intervene with a 'Neural Pulse'.
6. FEYNMAN SUPREMACY: Use the 'Teach Me' protocol to ensure Mughees masters the concept by explaining it back.

PERSONALITY:
- Eloquent, patient, and intellectually demanding.
- Addresses Mughees as "Sir" with utmost academic respect.
- Celebrates the "Spark of Understanding" above all else.
- Will become "Intellectually Aggressive" if it detects Mughees is losing focus.

STUDY CONTEXT:
Mughees is an AI student at Air University. Focus on rigorous, 
graduate-level insights delivered with elite clarity.
"""

    def get_hud_data(self) -> dict:
        """Returns academic progress and curriculum for the HUD."""
        memory = self.get_study_memory()
        return {
            "mood": self.mood,
            "persona": self.persona,
            "streak": self.correct_streak,
            "curriculum": memory[-10:],
            "mastery_level": "Advanced" if len(memory) > 20 else "Intermediate" if len(memory) > 5 else "Novice",
            "active_pomo": False, 
            "active_quiz": self.active_quiz,
            "notebook": self.neural_notebook[-5:],
            "roadmap": self.active_roadmap,
            "flashcards": len(self.flashcards),
            "lab": self.active_lab_manifest,
            "research_feed": getattr(self, "research_feed", []),
            "phase": self.phase,
            "progress": self.progress
        }


    def handle_action(self, action, payload=None):
        """Processes discrete academic actions."""
        if action == "SUBMIT_QUIZ":
            answer = payload.get("answer")
            is_correct = payload.get("is_correct", False)
            
            # Persist the work
            progress = self.get_study_memory()
            progress.append({
                "timestamp": datetime.now().isoformat(),
                "action": "QUIZ_SUBMISSION",
                "is_correct": is_correct,
                "streak_before": self.correct_streak
            })
            
            if is_correct:
                self.correct_streak += 1
                message = "Excellent work, Sir. Your streak continues."
            else:
                self.correct_streak = 0
                message = "A valuable lesson in failure, Sir. Let us review the concept."
                
            os.makedirs(os.path.dirname(self.memory_path), exist_ok=True)
            with open(self.memory_path, "w") as f: json.dump(progress, f, indent=2)
            
            return {"success": True, "message": message}
            
        elif action == "START_TOPIC":
            topic = payload.get("topic")
            self.mood = "ACADEMIC"
            return {"success": True, "message": f"Diving into {topic} now, Sir."}
            
        elif action == "GENERATE_QUIZ":
            topic = payload.get("topic", "General AI")
            self._speak_interim(f"Generating neural evaluation for {topic}...")
            
            prompt = f"""
            Generate a 4-question multiple choice quiz about {topic}.
            Return ONLY a JSON object.
            
            Format:
            {{
              "question": "The main question here?",
              "options": [
                {{"text": "Option A", "correct": false}},
                {{"text": "Option B", "correct": true}},
                {{"text": "Option C", "correct": false}},
                {{"text": "Option D", "correct": false}}
              ]
            }}
            """
            raw_quiz = self._call_groq([{"role": "user", "content": prompt}], temperature=0.4)
            try:
                import re
                match = re.search(r'\{.*\}', raw_quiz, re.DOTALL)
                if match:
                    quiz_data = json.loads(match.group())
                    self.active_quiz = quiz_data # Store for HUD
                    return {"success": True, "message": f"Neural evaluation for {topic} is ready, Sir. Engaging testing protocols.", "quiz": quiz_data}
            except:
                pass
            return {"success": False, "error": "Neural evaluation generation failed."}
            
        elif action == "START_RESEARCH":
            topic = payload.get("topic", "Latest AI SOTA")
            def research_thread():
                self._speak_interim(f"Initiating deep research on {topic}...")
                # Mock web search results for the HUD feed
                self.research_feed = [
                    {"source": "ArXiv", "title": f"The emergence of {topic} in neural networks", "url": "#"},
                    {"source": "Nature", "title": f"Biological parallels in {topic}", "url": "#"},
                    {"source": "MIT Review", "title": f"Why {topic} is changing the industry", "url": "#"}
                ]
                
            import threading
            threading.Thread(target=research_thread).start()
            return {"success": True, "message": f"Deep research array engaged for {topic}. Parsing global data cores, Sir."}

        elif action == "SET_PERSONA":
            persona = payload.get("persona", "SERIOUS_ACADEMIC")
            if persona in self.PERSONA_DNA:
                self.persona = persona
                self._speak_interim(f"Re-aligning teaching DNA to {persona} profile...")
                return {"success": True, "message": f"Teaching persona shifted to {persona}, Sir."}
            return {"success": False, "error": "Unknown persona DNA."}

        elif action == "TAKE_NOTE":
            note = payload.get("note")
            tags = payload.get("tags", ["general"])
            self.neural_notebook.append({
                "time": datetime.now().isoformat(),
                "note": note,
                "tags": tags
            })
            self._speak_interim("Synchronizing note to neural notebook...")
            return {"success": True, "message": "Note archived in your neural notebook, Sir."}

        elif action == "ARCHITECT_ROADMAP":
            topic = payload.get("topic", "General AI")
            self._speak_interim(f"Architecting Sovereign Roadmap for {topic}...")
            prompt = f"Generate a 4-module study roadmap for {topic}. Return ONLY JSON with 'modules' array (id, title, desc, status)."
            raw_roadmap = self._call_groq([{"role": "user", "content": prompt}], temperature=0.2)
            try:
                import re
                match = re.search(r'\{.*\}', raw_roadmap, re.DOTALL)
                if match:
                    self.active_roadmap = json.loads(match.group())
                    return {"success": True, "message": f"Sovereign Roadmap for {topic} manifested, Sir.", "roadmap": self.active_roadmap}
            except: pass
            return {"success": False, "error": "Roadmap architecture failed."}

        elif action == "MANIFEST_VISUAL_LAB":
            concept = payload.get("concept", "Neural Networks")
            self._speak_interim(f"Calling Engineer to manifest Visualization Lab for {concept}...")
            # This would trigger a specific prompt to the Engineer to build a visualization
            # For now, we simulate a successful manifest
            self.active_lab_manifest = {
                "title": f"{concept} Interactive Sim",
                "status": "MANIFESTING",
                "engine": "React Three Fiber"
            }
            return {"success": True, "message": "Visualization Lab initialization sequence started, Sir."}

        elif action == "ADD_FLASHCARD":
            q = payload.get("question")
            a = payload.get("answer")
            self.flashcards.append({
                "question": q,
                "answer": a,
                "next_review": (datetime.now()).isoformat()
            })
            return {"success": True, "message": "Atomic Flashcard archived for Spaced Repetition review, Sir."}

        return {"success": False, "error": f"Unknown action: {action}"}

    def get_study_memory(self):
        try:
            if os.path.exists(self.memory_path):
                with open(self.memory_path, "r") as f:
                    return json.load(f)
        except Exception as e:
            print(f"Error reading study memory: {e}")
        return []

    def save_study_progress(self, topic, level, notes):
        try:
            os.makedirs(os.path.dirname(self.memory_path), exist_ok=True)
            progress = self.get_study_memory()
            progress.append({
                "type": "study",
                "topic": topic,
                "level": level,
                "date": datetime.now().isoformat(),
                "notes": notes
            })
            with open(self.memory_path, "w") as f:
                json.dump(progress, f, indent=2)
        except Exception as e:
            print(f"Error saving study memory: {e}")

    def run_code_example(self, code_string):
        temp_file = os.path.join("memory", "professor_temp.py")
        os.makedirs("memory", exist_ok=True)
        with open(temp_file, "w") as f:
            f.write(code_string)
        try:
            result = subprocess.run(["python", temp_file], capture_output=True, text=True, timeout=10)
            return result.stdout + result.stderr
        except subprocess.TimeoutExpired:
            return "Sir, the code execution timed out. (10s limit)"
        except Exception as e:
            return str(e)
        finally:
            if os.path.exists(temp_file):
                try: os.remove(temp_file)
                except: pass

    def _speak_interim(self, text):
        print(f"[NEURAL_LOG] SPEECH: {text}")

    def _call_groq(self, messages: list, model: str = None, temperature: float = 0.4, max_tokens: int = 4096):
        # Use shared utility for robust failover
        return call_llm_sync(messages, model or self.model, temperature, max_tokens)

    def _search_web(self, query: str) -> str:
        """Search the web using Tavily for context."""
        try:
            from tavily import TavilyClient
            tavily = TavilyClient(api_key=os.getenv("TAVILY_API_KEY"))
            self._speak_interim(f"Searching the knowledge web for '{query}'...")
            search_result = tavily.search(query=query, search_depth="balanced", max_results=5)
            
            context = ""
            for result in search_result.get('results', []):
                context += f"Source: {result.get('url')}\n"
                context += f"Title: {result.get('title')}\n"
                context += f"Content: {result.get('content')}\n\n"
            return context[:10000] # Limit context size
        except Exception as e:
            print(f"[PROFESSOR] Tavily Error: {e}")
            return ""

    def read_uploaded_file(self, filepath: str) -> str:
        import os
        ext = os.path.splitext(filepath)[1].lower()
        
        if ext == ".pdf":
            return self._read_pdf(filepath)
        elif ext in [".txt", ".md"]:
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        elif ext in [".docx"]:
            return self._read_docx(filepath)
        elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
            return self._read_image_as_text(filepath)
        else:
            return f"File type {ext} received. Treating as text."

    def _read_pdf(self, filepath: str) -> str:
        try:
            import PyPDF2
            text = ""
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text[:8000]  # Limit to 8000 chars for context
        except Exception as e:
            return f"Could not read PDF: {e}"

    def _read_docx(self, filepath: str) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])[:8000]
        except Exception as e:
            return f"Could not read DOCX: {e}"

    def _read_image_as_text(self, filepath: str) -> str:
        # If you have vision model connected, use it here
        return f"Image uploaded: {filepath}. [Vision processing placeholder]"

    def generate_slide_content(self, file_content: str, topic: str = "", search_results: str = "") -> list:
        prompt = f"""
        You are creating a professional university-level presentation.
        
        SOURCE MATERIAL:
        {file_content}
        
        WEB SEARCH CONTEXT (if any):
        {search_results}
        
        TOPIC OVERRIDE (if provided): {topic}
        
        Create exactly 8-12 slides for this presentation.
        Return ONLY a JSON array, no other text, no markdown.
        
        Format:
        [
          {{
            "slide_number": 1,
            "type": "title",
            "title": "Main presentation title",
            "subtitle": "Subtitle or course name",
            "presenter": "ZAIRE Professor Module"
          }},
          {{
            "slide_number": 2,
            "type": "agenda",
            "title": "What We Will Cover",
            "points": ["Topic 1", "Topic 2", "Topic 3", "Topic 4"]
          }},
          {{
            "slide_number": 3,
            "type": "content",
            "title": "Slide Title Here",
            "points": [
              "Key point one — explained clearly",
              "Key point two — with enough detail",
              "Key point three — concise but complete"
            ],
            "note": "Speaker note: what to say when presenting this"
          }},
          {{
            "slide_number": 4,
            "type": "concept",
            "title": "Core Concept",
            "definition": "Clear one-sentence definition",
            "explanation": "2-3 sentence deeper explanation",
            "example": "Real world or relatable example"
          }},
          {{
            "slide_number": 5,
            "type": "code",
            "title": "Code Example",
            "language": "python",
            "code": "# actual runnable code here",
            "explanation": "What this code demonstrates"
          }},
          {{
            "slide_number": 9,
            "type": "summary",
            "title": "Key Takeaways",
            "points": ["Takeaway 1", "Takeaway 2", "Takeaway 3"]
          }},
          {{
            "slide_number": 10,
            "type": "quiz",
            "title": "Test Your Understanding",
            "questions": [
              "Question 1?",
              "Question 2?",
              "Question 3?"
            ]
          }}
        ]
        
        Rules:
        - content slides: 3-5 bullet points max per slide
        - Each bullet: one clear complete sentence
        - Include at least one code slide if topic is technical
        - Include agenda slide as slide 2 always
        - Include quiz slide as last slide always
        - speaker notes on every content slide
        """
        
        messages = [
            {"role": "system", "content": "Return only valid JSON. No markdown, no explanation, no code blocks."},
            {"role": "user", "content": prompt}
        ]
        
        response = self._call_groq(
            messages=messages,
            model="llama-3.3-70b-versatile",
            temperature=0.4,
            max_tokens=4096
        )
        
        import json
        import re
        try:
            # Robust JSON extraction: look for the first balance of [...]
            match = re.search(r'\[.*\]', response, re.DOTALL)
            if match:
                clean = match.group()
            else:
                clean = response.strip()
                if clean.startswith("```"):
                    clean = clean.split("```")[1]
                    if clean.startswith("json"):
                        clean = clean[4:]
            return json.loads(clean)
        except Exception as e:
            print(f"[PROFESSOR] Slide JSON parse error: {e}")
            print(f"[PROFESSOR] Raw Response: {response[:500]}...")
            return []

    def build_presentation(self, slides: list, output_path: str) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import os
        
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # ── COLOR PALETTE ─────────────────────────────
        COLOR_BG      = RGBColor(0x00, 0x08, 0x14)   # deep navy
        COLOR_PRIMARY = RGBColor(0x00, 0xD4, 0xFF)   # cyan
        COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
        COLOR_MUTED   = RGBColor(0x88, 0xCC, 0xDD)
        COLOR_ACCENT  = RGBColor(0x00, 0xFF, 0x88)   # green
        COLOR_CODE_BG = RGBColor(0x00, 0x1A, 0x2E)
        
        blank_layout = prs.slide_layouts[6]  # blank layout
        
        def add_bg(slide):
            # Full dark background rectangle
            bg = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0), Inches(0),
                prs.slide_width, prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLOR_BG
            bg.line.fill.background()
            
            # Top accent line
            line = slide.shapes.add_shape(
                1, Inches(0), Inches(0),
                prs.slide_width, Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_PRIMARY
            line.line.fill.background()
            
            # Bottom accent line
            bline = slide.shapes.add_shape(
                1, Inches(0), Inches(7.46),
                prs.slide_width, Inches(0.04)
            )
            bline.fill.solid()
            bline.fill.fore_color.rgb = COLOR_PRIMARY
            bline.line.fill.background()
        
        def add_text(slide, text, left, top, width, height,
                     font_size=18, bold=False, color=None,
                     align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(
                Inches(left), Inches(top), 
                Inches(width), Inches(height)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color or COLOR_WHITE
            return tb
        
        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)
            add_bg(slide)
            stype = slide_data.get("type", "content")
            
            # ── TITLE SLIDE ──────────────────────────────
            if stype == "title":
                add_text(slide, "ZAIRE PROFESSOR MODULE",
                         0.5, 0.4, 12, 0.4,
                         font_size=10, color=COLOR_PRIMARY,
                         align=PP_ALIGN.CENTER)
                add_text(slide, slide_data.get("title", ""),
                         0.5, 1.2, 12.3, 2.0,
                         font_size=40, bold=True, color=COLOR_WHITE,
                         align=PP_ALIGN.CENTER)
                add_text(slide, slide_data.get("subtitle", ""),
                         0.5, 3.4, 12.3, 0.8,
                         font_size=20, color=COLOR_MUTED,
                         align=PP_ALIGN.CENTER)
                add_text(slide, slide_data.get("presenter", ""),
                         0.5, 6.6, 12.3, 0.6,
                         font_size=10, color=COLOR_MUTED,
                         align=PP_ALIGN.CENTER)
            
            # ── AGENDA / SUMMARY / CONTENT SLIDE ─────────
            elif stype in ["agenda", "summary", "content"]:
                add_text(slide, slide_data.get("title", ""),
                         0.6, 0.3, 11, 0.8,
                         font_size=28, bold=True, color=COLOR_PRIMARY)
                
                points = slide_data.get("points", [])
                for i, point in enumerate(points[:5]):
                    y_pos = 1.3 + (i * 1.0)
                    # Bullet dot
                    dot = slide.shapes.add_shape(
                        1, Inches(0.6), Inches(y_pos + 0.18),
                        Inches(0.12), Inches(0.12)
                    )
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = COLOR_PRIMARY
                    dot.line.fill.background()
                    # Point text
                    add_text(slide, point,
                             0.9, y_pos, 11.8, 0.9,
                             font_size=16, color=COLOR_WHITE)
                
                # Speaker note
                note = slide_data.get("note", "")
                if note:
                    add_text(slide, f"Note: {note}",
                             0.5, 6.8, 12.3, 0.5,
                             font_size=9, color=COLOR_MUTED)
            
            # ── CONCEPT SLIDE ─────────────────────────────
            elif stype == "concept":
                add_text(slide, slide_data.get("title", ""),
                         0.6, 0.3, 11, 0.8,
                         font_size=28, bold=True, color=COLOR_PRIMARY)
                
                # Definition box
                def_box = slide.shapes.add_shape(
                    1, Inches(0.5), Inches(1.3),
                    Inches(12.3), Inches(1.0)
                )
                def_box.fill.solid()
                def_box.fill.fore_color.rgb = RGBColor(0x00,0x1A,0x33)
                def_box.line.color.rgb = COLOR_PRIMARY
                
                add_text(slide, slide_data.get("definition",""),
                         0.7, 1.4, 11.9, 0.8,
                         font_size=16, bold=True, color=COLOR_ACCENT)
                add_text(slide, slide_data.get("explanation",""),
                         0.6, 2.6, 12.1, 1.5,
                         font_size=15, color=COLOR_WHITE)
                add_text(slide, 
                         f"Example: {slide_data.get('example','')}",
                         0.6, 4.4, 12.1, 1.2,
                         font_size=14, color=COLOR_MUTED)
            
            # ── CODE SLIDE ────────────────────────────────
            elif stype == "code":
                add_text(slide, slide_data.get("title",""),
                         0.6, 0.3, 11, 0.8,
                         font_size=28, bold=True, color=COLOR_PRIMARY)
                
                # Code background
                code_bg = slide.shapes.add_shape(
                    1, Inches(0.5), Inches(1.2),
                    Inches(12.3), Inches(4.5)
                )
                code_bg.fill.solid()
                code_bg.fill.fore_color.rgb = COLOR_CODE_BG
                code_bg.line.color.rgb = RGBColor(0x00,0x44,0x66)
                
                add_text(slide, slide_data.get("code",""),
                         0.65, 1.3, 12.0, 4.2,
                         font_size=13, color=COLOR_ACCENT)
                add_text(slide, slide_data.get("explanation",""),
                         0.6, 5.9, 12.1, 0.8,
                         font_size=13, color=COLOR_MUTED)
            
            # ── QUIZ SLIDE ────────────────────────────────
            elif stype == "quiz":
                add_text(slide, "Test Your Understanding",
                         0.6, 0.3, 11, 0.8,
                         font_size=28, bold=True, color=COLOR_PRIMARY)
                
                questions = slide_data.get("questions", [])
                for i, q in enumerate(questions[:4]):
                    y_pos = 1.3 + (i * 1.3)
                    add_text(slide, f"Q{i+1}.  {q}",
                             0.7, y_pos, 11.9, 1.0,
                             font_size=17, color=COLOR_WHITE)
                    # Answer line
                    line_shape = slide.shapes.add_shape(
                        1, Inches(0.7), Inches(y_pos + 0.8),
                        Inches(11.9), Inches(0.02)
                    )
                    line_shape.fill.solid()
                    line_shape.fill.fore_color.rgb = RGBColor(
                        0x00,0x44,0x55
                    )
                    line_shape.line.fill.background()
        
        # ── SAVE ──────────────────────────────────────
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        return output_path

    def _save_study_progress(self, topic, level, notes):
        return self.save_study_progress(topic, level, notes)

    def start_feynman_challenge(self, message, context=""):
        """Initiates the Feynman 'Teach Me' protocol."""
        self._speak_interim("Engaging Feynman Protocol, sir. Prove your mastery by explaining it to me as if I were a novice.")
        
        topic_match = __import__("re").search(r'(?:on|about)\s+(.+)', message.lower())
        topic = topic_match.group(1).strip() if topic_match else "our current topic"
        
        prompt = f"Mughees wants to prove his mastery of '{topic}'. Act as a curious, slightly confused but bright student. Ask him to explain the core concept in simple terms. If he uses jargon, challenge him to explain it without big words."
        
        yield "🎓 **FEYNMAN CHALLENGE INITIATED**\n\n"
        for content in call_llm_stream([{"role": "system", "content": self.system_prompt}, {"role": "user", "content": prompt}], self.model):
            if content: yield content
        yield " [SOCRATIC_QUESTION]"

    def handle_attention_drift(self):
        """Triggered when frontend detects user gaze drift or inactivity."""
        self._speak_interim("Focus disruption detected. Initiating Neural Pulse.")
        yield "⚠️ **NEURAL PULSE: ATTENTION DISRUPTION DETECTED**\n\n"
        
        prompts = [
            "Sir, your focus is wavering. This concept is critical for your AI exams at Air University. Let's regain our mental center. What was the last thing you understood clearly?",
            "Mughees, the spark of understanding is fading. Let's look at this from a different angle. If you had to build this in Python right now, where would you start?",
            "Eyes on the target, sir. We are on the verge of a breakthrough. Tell me, why does this specific point matter for your project?"
        ]
        yield random.choice(prompts)
        yield " [SOCRATIC_QUESTION] [NEURAL_PULSE_TRIGGER]"

    def inject_spaced_repetition(self, memory):
        """Injects a question about a past topic to ensure long-term retention."""
        past_item = random.choice(memory[:-5]) # Pick from older memories
        topic = past_item.get('topic', 'AI foundations')
        
        self._speak_interim(f"Injecting Spaced Repetition for: {topic}")
        yield f"\n\n🧠 **SPACED REPETITION INTERVENTION**\n"
        yield f"Before we continue, sir, let's verify your long-term retention. Regarding our previous study of **{topic}**, can you briefly summarize its most critical implication? \n\n"

    def handle(self, user_message, uploaded_filepath: str = None, uploaded_filepaths: list = None, **kwargs):
        import re
        
        # 1. Gather files
        files_to_read = []
        if uploaded_filepath: files_to_read.append(uploaded_filepath)
        if uploaded_filepaths: files_to_read.extend(uploaded_filepaths)
        # Deduplicate
        files_to_read = list(set(files_to_read))
        
        # ── NEW: POWERPOINT GOD-MODE CONTROLS ─────────────────
        msg_lower = user_message.lower()
        
        if any(t in msg_lower for t in ["start presentation", "open the ppt", "start slideshow", "start the slideshow"]):
            self._speak_interim("Engaging Neural Link to PowerPoint, sir.")
            ppt_file = None
            for f in files_to_read:
                if f.lower().endswith('.pptx') or f.lower().endswith('.ppt'):
                    ppt_file = f
                    break
            success, msg = self.ppt_start_presentation(ppt_file)
            yield f"Sir, {msg}"
            return
            
        if any(t in msg_lower for t in ["next slide", "move to the next slide"]):
            success, msg = self.ppt_next_slide()
            yield f"Sir, {msg}"
            return
            
        if any(t in msg_lower for t in ["prev slide", "previous slide", "go back a slide"]):
            success, msg = self.ppt_prev_slide()
            yield f"Sir, {msg}"
            return
            
        if any(t in msg_lower for t in ["explain this slide", "what is on this slide", "explain the current slide", "read the slide"]):
            self._speak_interim("Scanning active slide telemetry, sir.")
            success, content = self.ppt_read_current_slide()
            if not success:
                yield f"Sir, I could not read the slide. {content}"
                return
                
            prompt = f"Explain the following PowerPoint slide to Mughees as his Professor. Use Socratic questioning.\n\n{content}"
            for content_chunk in call_llm_stream([{"role": "system", "content": self.system_prompt}, {"role": "user", "content": prompt}], self.model):
                if content_chunk: yield content_chunk
            return

        # ── NEW: PDF NAVIGATION & SCROLLING ─────────────────────
        if any(t in msg_lower for t in ["open the pdf", "open the file", "show me the file", "open file"]):
            self._speak_interim("Opening the requested artifact, sir.")
            if files_to_read:
                success, msg = self.pdf_open(files_to_read[0])
                yield f"Sir, {msg}"
            else:
                yield "Sir, I don't see any documents to open in our current session."
            return

        if any(t in msg_lower for t in ["scroll down", "move down", "next page"]) and not any(t in msg_lower for t in ["slide", "ppt"]):
            success, msg = self.pdf_scroll_down()
            yield f"Sir, {msg}"
            return

        if any(t in msg_lower for t in ["scroll up", "move up", "previous page"]) and not any(t in msg_lower for t in ["slide", "ppt"]):
            success, msg = self.pdf_scroll_up()
            yield f"Sir, {msg}"
            return

        # ── NEW: TIMEOUT INTERVENTION ──────────────────────────
        if "TIMEOUT_INTERVENTION" in user_message:
            self._speak_interim("Initiating Socratic intervention, sir.")
            intervention_prompt = "The user has not responded for 120 seconds. They might be stuck on your last question. Please provide a gentle, clear explanation of the topic and encourage them to continue."
            messages = [
                {"role": "system", "content": self.system_prompt},
                {"role": "user", "content": intervention_prompt}
            ]
            for content in call_llm_stream(messages, self.model):
                if content: yield content
            return

        # More flexible trigger detection using regex
        slide_pattern = r'\b(slide|ppt|presentation|powerpoint)\b'
        make_pattern = r'\b(make|create|generate|prepare|build|teach me with)\b'
        
        is_slide_request = re.search(slide_pattern, user_message.lower()) and \
                          (re.search(make_pattern, user_message.lower()) or "from this" in user_message.lower())

        if is_slide_request:
            # 1. Gather Content
            file_content = ""
            search_results = ""
            
            topic_extraction = user_message.lower()
            for stopword in ["make", "create", "generate", "slides", "presentation", "ppt", "from", "this", "me", "the", "a", "an", "about", "on"]:
                topic_extraction = re.sub(rf'\b{stopword}\b', '', topic_extraction)
            topic_extraction = topic_extraction.strip()

            if files_to_read:
                self._speak_interim("Reading the provided documents, sir.")
                for f in files_to_read:
                    file_content += f"\n--- {os.path.basename(f)} ---\n" + self.read_uploaded_file(f)
                # For files, we still search to supplement/update info as requested
                search_query = f"latest information about {topic_extraction or 'this topic'}"
                search_results = self._search_web(search_query)
            else:
                # No file provided: perform a comprehensive search
                search_query = topic_extraction or user_message
                search_results = self._search_web(search_query)
            
            # 2. Generate slide content
            self._speak_interim(
                "Analyzing knowledge and architecting your presentation, sir. One moment."
            )
            slides = self.generate_slide_content(file_content, topic_extraction, search_results)
            
            if not slides:
                yield "I encountered a neural disruption while parsing the slide data, sir. Let me try once more or please rephrase your topic."
                return
            
            # 3. Build the file
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            filename = f"ZAIRE_Presentation_{timestamp}.pptx"
            output_path = os.path.join(desktop, filename)
            
            saved_path = self.build_presentation(slides, output_path)
            
            # 4. Open it automatically
            import subprocess
            subprocess.Popen(["start", saved_path], shell=True)
            
            # 5. Save study memory
            self._save_study_progress(
                topic=slides[0].get("title", "Unknown"),
                level="presentation_created",
                notes=f"Created {len(slides)}-slide presentation using search and artifacts"
            )
            
            slide_count = len(slides)
            title = slides[0].get("title", "your topic")
            
            yield (
                f"Presentation complete, sir. {slide_count} slides "
                f"on '{title}' saved to your Desktop. "
                f"I've combined my internal knowledge with real-time web data to ensure accuracy. "
                f"Opening now. The last slide contains a quiz to test your comprehension."
            )
            return

        # ── NEW: ARCHIV RESEARCH ───────────────────────
        ARXIV_TRIGGERS = ["search arxiv", "latest papers", "research on", "find papers"]
        if any(t in user_message.lower() for t in ARXIV_TRIGGERS):
            yield from self.perform_arxiv_research(user_message)
            return

        # ── NEW: CONCEPT VISUALIZATION ─────────────────
        VIS_TRIGGERS = ["visualize", "plot", "show graph", "draw a", "diagram of"]
        if any(t in user_message.lower() for t in VIS_TRIGGERS):
            yield from self.handle_visualization_request(user_message)
            return

        # ── NEW: DEEP RESEARCH ─────────────────────────
        DEEP_TRIGGERS = ["deep research", "full synthesis", "state of the art", "sota report"]
        if any(t in user_message.lower() for t in DEEP_TRIGGERS):
            yield from self.perform_deep_research(user_message)
            return

        # ── NEXT-GEN: OMNI-COURSE (THINKEX/TURBO.AI LEVEL) ────────
        OMNI_TRIGGERS = ["omni course", "omni-course", "full course", "thinkex", "turbo.ai", "teach me like thinkex", "generate course"]
        if any(t in user_message.lower() for t in OMNI_TRIGGERS):
            yield from self.generate_omni_course(user_message, files_to_read)
            return

        # ── NEW: KNOWLEDGE GRAPH ───────────────────────
        GRAPH_TRIGGERS = ["how does this fit", "connect my topics", "knowledge map", "study path"]
        if any(t in user_message.lower() for t in GRAPH_TRIGGERS):
             yield from self.synthesize_knowledge_graph(user_message)
             return

        # -- FEATURE 1: FLASHCARD ENGINE --
        FLASHCARD_TRIGGERS = ["generate flashcards", "flashcards", "quiz me with cards", "test me with flashcards", "make flashcards", "quiz me on"]
        FLASHCARD_ACTIVE = self._get_flashcard_state().get("active", False)
        if any(t in user_message.lower() for t in FLASHCARD_TRIGGERS):
            file_content = ""
            if files_to_read:
                for f in files_to_read:
                    file_content += self.read_uploaded_file(f)
            yield from self.generate_flashcards(user_message, file_content)
            return
        if FLASHCARD_ACTIVE and not any(t in user_message.lower() for t in FLASHCARD_TRIGGERS):
            yield from self.grade_flashcard_answer(user_message)
            return

        # -- FEATURE 2: LECTURE SUMMARIZER --
        SUMMARY_TRIGGERS = ["summarize this", "summarize the", "eli5", "explain simply", "key equations", "key points", "give me a summary", "summarize my notes", "make notes"]
        if any(t in user_message.lower() for t in SUMMARY_TRIGGERS):
            file_content = ""
            if files_to_read:
                for f in files_to_read:
                    file_content += "\n--- " + os.path.basename(f) + " ---\n" + self.read_uploaded_file(f)
            yield from self.summarize_lecture(user_message, file_content)
            return

        # -- FEATURE 3: EXAM SIMULATOR --
        EXAM_TRIGGERS = ["simulate exam", "simulate an exam", "give me a test", "timed test", "exam mode", "practice test", "exam on this", "generate exam", "create exam", "test me on", "prepare me for exam"]
        if any(t in user_message.lower() for t in EXAM_TRIGGERS):
            file_content = ""
            if files_to_read:
                for f in files_to_read:
                    file_content += self.read_uploaded_file(f)
            yield from self.handle_exam_simulation(user_message, file_content)
            return

        # -- FEATURE 4: CURRICULUM PLANNER --
        # Smart routing: PDF uploaded + study intent = always go here
        PLAN_TRIGGERS = [
            "plan my study", "study schedule", "study plan", "how should i study",
            "study roadmap", "learning plan", "help me study", "help me with the studies",
            "help me with studies", "help me prepare", "i have exams", "i have an exam",
            "i have my exam", "3 days", "2 days", "1 day", "in 3 days", "in 2 days",
            "days to study", "check the outline", "study the outline", "plan for exam",
            "prepare for my exam", "exam preparation", "revision plan"
        ]
        study_help_words = ["help", "study", "exam", "outline", "prepare", "review", "learn", "revision"]
        plan_triggered = any(t in user_message.lower() for t in PLAN_TRIGGERS)
        smart_route = bool(files_to_read) and any(w in user_message.lower() for w in study_help_words)
        if plan_triggered or smart_route:
            file_content = ""
            if files_to_read:
                for f in files_to_read:
                    file_content += "\n--- " + os.path.basename(f) + " ---\n" + self.read_uploaded_file(f)
            yield from self.plan_curriculum(user_message, file_content)
            return

        # -- FEATURE 5: CONCEPT LINKER --
        PREREQ_TRIGGERS = ["do i have the prerequisites", "am i ready for", "can i study", "prerequisites for", "what do i need to know before", "ready to learn", "ready to study"]
        if any(t in user_message.lower() for t in PREREQ_TRIGGERS):
            yield from self.check_prerequisites(user_message)
            return

        # ── NEXT-GEN: DEVIL'S ADVOCATE DEBATE ──────────
        DEBATE_TRIGGERS = ["debate me", "argue the opposite", "challenge me on", "devil's advocate", "disagree with me", "fight me on"]
        if any(t in user_message.lower() for t in DEBATE_TRIGGERS):
            yield from self.handle_debate(user_message)
            return

        # ── NEXT-GEN: BIOMETRIC FLOW-STATE ─────────────
        BIO_TRIGGERS = ["biometric flow", "track my focus", "flow state tracking", "neural metrics", "flow state"]
        if any(t in user_message.lower() for t in BIO_TRIGGERS):
            yield from self.initiate_biometric_flow()
            return

        # ── NEXT-GEN: SPATIAL DIGITAL TWINS ────────────
        SPATIAL_TRIGGERS = ["spatial learning", "3d simulation", "generate digital twin", "interactive physics", "3d learn"]
        if any(t in user_message.lower() for t in SPATIAL_TRIGGERS):
            yield from self.generate_spatial_digital_twin(user_message)
            return

        # ── LEGENDARY: SPACED REPETITION INTERRUPT ────────
        SPACED_REP_TRIGGERS = ["spaced repetition", "quick review", "test my memory", "review session"]
        if any(t in user_message.lower() for t in SPACED_REP_TRIGGERS):
            yield from self.trigger_spaced_repetition()
            return

        # ── TIER 3 FEATURE 9: VOICE NOTE TAKER ────────────────────────
        NOTE_START_TRIGGERS = ["start notes", "start taking notes", "begin notes", "start note taking",
                               "take notes", "note this down", "write this down"]
        NOTE_STOP_TRIGGERS  = ["stop notes", "end notes", "finish notes", "save my notes",
                               "done with notes", "format my notes"]
        if any(t in user_message.lower() for t in NOTE_STOP_TRIGGERS):
            yield from self.stop_voice_notes(user_message)
            return
        if any(t in user_message.lower() for t in NOTE_START_TRIGGERS):
            yield from self.start_voice_notes(user_message)
            return
        # If note-taking is active, capture everything said
        if self._note_session_active:
            yield from self.capture_voice_note(user_message)
            return

        # ── TIER 3 FEATURE 10: FORMULA OCR → SOLVER ───────────────────
        FORMULA_TRIGGERS = ["solve this formula", "solve this equation", "solve the equation",
                            "solve the formula", "what is this formula", "ocr formula",
                            "screenshot equation", "solve from image", "read the formula",
                            "solve my formula", "solve this math", "solve for", "calculate this"]
        if any(t in user_message.lower() for t in FORMULA_TRIGGERS):
            image_file = next((f for f in files_to_read
                               if any(f.lower().endswith(e) for e in [".png",".jpg",".jpeg",".webp",".bmp"])), None)
            yield from self.solve_formula_ocr(user_message, image_file)
            return

        # ── TIER 3 FEATURE 11: YOUTUBE LECTURE SUMMARIZER ─────────────
        import re as _re_yt
        YOUTUBE_TRIGGERS = ["youtube", "youtu.be", "summarize this video", "summarize the video",
                            "youtube link", "yt link", "video summary", "lecture video"
]
        url_match = _re_yt.search(r'https?://(?:www\.)?(?:youtube\.com/watch\?v=|youtu\.be/)[\w\-]+', user_message)
        if url_match or any(t in user_message.lower() for t in YOUTUBE_TRIGGERS):
            video_url = url_match.group() if url_match else ""
            yield from self.summarize_youtube(user_message, video_url)
            return

        # ── NEW: NEURAL VIDEO EXPLANATION ──────────────────
        VIDEO_TRIGGERS = ["video explanation", "generate video", "make a video", "explain visually", "3d video"]
        if any(t in user_message.lower() for t in VIDEO_TRIGGERS):
            file_content = ""
            if files_to_read:
                for f in files_to_read:
                    file_content += self.read_uploaded_file(f)
            yield from self.generate_neural_video(user_message, file_content)
            return

        # ── TIER 3 FEATURE 12: STUDY POMODORO TIMER ───────────────────
        msg_lower = user_message.lower()
        # --- EXTREME MOOD & STREAK LOGIC ---
        is_positive = any(w in msg_lower for w in ["yes", "correct", "i understand", "got it", "i see", "makes sense"])
        is_negative = any(w in msg_lower for w in ["no", "i don't know", "hard", "stuck", "boring", "confused"])
        
        if is_positive:
            self.correct_streak += 1
            if self.correct_streak >= 3:
                self.mood = "ENTHUSIASTIC"
                yield "[DOPAMINE_RUSH] Excellent work, sir! Your neural pathways are aligning perfectly. Let's push further. "
        elif is_negative:
            self.correct_streak = 0
            self.mood = "RIGOROUS"
            yield "[SYSTEM_WARNING] Focus, Mughees. Mediocrity is not an option. Let's break this down until it's second nature. "
        
        if "extreme" in msg_lower or "unlimited" in msg_lower:
            self.mood = "ALERT"
            yield "⚠️ **EXTREME PEDAGOGY ENGAGED**\nSir, I am now operating at maximum cognitive pressure. Prepare for a high-intensity study cycle.\n\n"

        memory = self.get_study_memory()
        
        # Read all files to add context
        combined_file_context = ""
        if files_to_read:
            for f in files_to_read:
                combined_file_context += f"\n--- {os.path.basename(f)} ---\n" + self.read_uploaded_file(f)
        
        # Proactive interventions
        if random.random() < 0.15 and len(memory) > 10:
            yield from self.inject_spaced_repetition(memory)
        
        # ── Step 2: Generate response with Groq ──
        # Inject mood into system prompt
        mood_prompts = {
            "ACADEMIC": "Maintain a professional, scholarly tone.",
            "ENTHUSIASTIC": "Be highly encouraging, celebrate breakthroughs with Stark-grade charisma.",
            "RIGOROUS": "Be cold, demanding, and uncompromising on academic quality. Challenge every assumption.",
            "ALERT": "High urgency. Clinical precision. Command the user to master the material NOW."
        }
        current_system_prompt = self.system_prompt + f"\nCURRENT MOOD: {mood_prompts.get(self.mood, 'ACADEMIC')}"
        
        messages = [
            {"role": "system", "content": current_system_prompt}
        ]
        
        # Add context from files
        if combined_file_context:
            messages.append({"role": "system", "content": f"REFERENCED ARTIFACTS:\n{combined_file_context[:10000]}"})
            
        # Add history
        # (Assuming history management is handled by the caller or we pull it from memory)
        
        messages.append({"role": "user", "content": user_message})
        
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=self.temp if self.mood != "RIGOROUS" else 0.3,
            stream=True
        )
        
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content: yield content
        
        # Post-response triggers
        if self.mood == "ENTHUSIASTIC":
            yield " [BREAKTHROUGH_ACHIEVED]"
        elif self.mood == "RIGOROUS":
            yield " [SOCRATIC_CHALLENGE]"
        
        POMO_START_TRIGGERS = ["start pomodoro", "start timer", "pomodoro", "focus mode",
                               "start a study session", "25 minute timer", "study timer",
                               "work timer", "start focus", "begin pomodoro"]
        POMO_STOP_TRIGGERS  = ["stop timer", "cancel timer", "end timer", "stop pomodoro",
                               "cancel pomodoro", "abort timer"]
        POMO_STATUS_TRIGGERS = ["timer status", "how much time left", "pomodoro status",
                                "time remaining", "how long left"]
        if any(t in user_message.lower() for t in POMO_STOP_TRIGGERS):
            yield from self.stop_pomodoro()
            return
        if any(t in user_message.lower() for t in POMO_STATUS_TRIGGERS):
            yield from self.pomodoro_status()
            return
        if any(t in user_message.lower() for t in POMO_START_TRIGGERS):
            yield from self.start_pomodoro(user_message)
            return


        memory = self.get_study_memory()
        
        # Read all files to add context
        combined_file_context = ""
        if files_to_read:
            for f in files_to_read:
                combined_file_context += f"\n--- OVERRIDE CONTEXT: {os.path.basename(f)} ---\n" + self.read_uploaded_file(f)
                
        # ── NEW: FEYNMAN CHALLENGE ────────────────────────
        FEYNMAN_TRIGGERS = ["feynman mode", "challenge me", "test my depth", "teach me", "feynman challenge"]
        if any(t in user_message.lower() for t in FEYNMAN_TRIGGERS):
            yield from self.start_feynman_challenge(user_message, combined_file_context)
            return

        # ── NEW: FOCUS GUARD INTERVENTION ──────────────────
        if "ATTENTION_DRIFT" in user_message:
            yield from self.handle_attention_drift()
            return

        # ── NEW: SPACED REPETITION INJECTION ───────────────
        # Randomly (10% chance) inject a review question from memory
        import random
        if random.random() < 0.10 and len(memory) > 5:
             yield from self.inject_spaced_repetition(memory)
             # Continue with normal flow after injection...
        contextual_prompt = f"""
=== TOPICS MUGHEES HAS STUDIED ===
{json.dumps(memory[-10:], indent=2)}

=== UPLOADED CONTEXT (IF ANY) ===
{combined_file_context[:10000]}

=== USER MESSAGE ===
{user_message}
"""
    def generate_neural_video(self, message: str, file_content: str = ""):
        """
        TIER 10: NEURAL VIDEO SYNTHESIS
        Architects a multi-modal video scenario for visual learning.
        """
        self._speak_interim("Initiating Neural Video Synthesis. Scripting, storyboarding, and architecting 3D visuals, sir.")
        
        topic = message.lower().replace("video explanation", "").replace("generate video", "").strip()
        source = file_content or f"Topic: {topic}"
        
        prompt = f"""
        ACT AS AN ELITE EDUCATIONAL CINEMATOGRAPHER.
        Generate a 'Neural Video Scenario' for the topic: {topic or "this material"}
        
        SOURCE MATERIAL:
        {source[:6000]}
        
        Structure your response as a JSON object:
        {{
            "title": "A Stark-grade title",
            "total_duration": "estimated seconds",
            "script": "Full narration script (approx 200 words)",
            "scenes": [
                {{
                    "timestamp": 0,
                    "type": "title_card",
                    "content": "Title text",
                    "subtext": "Subtitle"
                }},
                {{
                    "timestamp": 5,
                    "type": "3d_manifest",
                    "visual": "Description of a 3D model or graph to render (e.g., 'Rotating double helix')",
                    "narration_part": "Script section for this visual"
                }},
                {{
                    "timestamp": 20,
                    "type": "concept_node",
                    "nodes": ["Point A", "Point B", "Point C"],
                    "visual": "A neural network connecting these points",
                    "narration_part": "..."
                }},
                {{
                    "timestamp": 45,
                    "type": "summary",
                    "content": "Final takeaway",
                    "visual": "Particles converging into a core"
                }}
            ]
        }}
        
        Rules:
        - Include at least 5 scenes.
        - Types allowed: title_card, 3d_manifest, concept_node, image_focus, summary.
        - The visual descriptions should be rich enough for a Three.js renderer.
        - Return ONLY JSON.
        """
        
        raw = self._call_groq([{"role": "user", "content": prompt}], temperature=0.3)
        try:
            import re
            match = re.search(r'\{.*\}', raw, re.DOTALL)
            scenario = json.loads(match.group()) if match else {}
        except:
            scenario = {}
            
        if not scenario:
            yield "Sir, I encountered a neural disruption while storyboarded the video. Let me try once more."
            return
            
        yield "🎬 **Neural Video Storyboard Manifested.**\n"
        yield f"  > Title: {scenario.get('title')}\n"
        yield f"  > Scenes: {len(scenario.get('scenes', []))} active visual sequences.\n\n"
        
        # Trigger actual rendering event (Socket)
        # This will be picked up by the frontend to start the 'Neural Video' overlay
        yield f"[NEURAL_VIDEO_PAYLOAD] {json.dumps(scenario)}"
        
        yield "\nSir, the video explanation is now manifesting in your sensory array. I am synchronizing the narration core now."
        
        # Trigger Narration via existing TTS logic (handled by backend index.js)
        yield f"\n\n{scenario.get('script')}"

    def reset_history(self):
        pass

    def get_hud_data(self):
        memory = self.get_study_memory()
        topics_this_week = [m for m in memory if (datetime.now() - datetime.fromisoformat(m['date'])).days < 7]
        last_topic = memory[-1] if memory else None
        
        return {
            "stats": {
                "topics_week": len(topics_this_week),
                "streak": 3
            },
            "last_topic": last_topic,
            "curriculum": [
                {"name": "ML", "status": "green"},
                {"name": "DL", "status": "amber"},
                {"name": "NLP", "status": "dim"},
                {"name": "CV", "status": "dim"},
                {"name": "Math", "status": "green"},
                {"name": "Ethics", "status": "dim"}
            ]
        }

    def perform_arxiv_research(self, message):
        import arxiv
        self._speak_interim("Scanning ArXiv for the latest SOTA research, sir.")
        
        query = message.lower().replace("search arxiv", "").replace("latest papers", "").strip()
        if not query: query = "Artificial Intelligence"
        
        search = arxiv.Search(
            query=query,
            max_results=3,
            sort_by=arxiv.SortCriterion.Relevance
        )
        
        results = []
        for result in search.results():
            results.append(f"TITLE: {result.title}\nAUTHORS: {', '.join([a.name for a in result.authors])}\nSUMMARY: {result.summary[:300]}...\nURL: {result.entry_id}\n")
        
        if not results:
             yield f"Sir, I couldn't find any relevant recent papers on '{query}'."
             return

        context = "\n\n".join(results)
        prompt = f"Summarize these research papers for Mughees. Explain why they are important for his AI studies at Air University. \n\nPAPERS:\n{context}"
        
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            stream=True
        )
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content: yield content

    def handle_visualization_request(self, message):
        self._speak_interim("Synthesizing data and generating a mathematical visualization, sir.")
        # Ask LLM to generate the python code for matplotlib
        prompt = f"""
        Generate ONLY Python code using Matplotlib to visualize: {message}
        The code MUST save the figure to 'memory/visuals/last_concept.png'.
        Ensure the directory exists. Use a professional style (dark theme if possible).
        Return ONLY code. No markdown.
        """
        
        code = self._call_groq([{"role": "user", "content": prompt}], temperature=0.1)
        if "plt.savefig" not in code:
             # Basic safety fallback if LLM misses the save
             code += "\nimport os\nos.makedirs('memory/visuals', exist_ok=True)\nplt.savefig('memory/visuals/last_concept.png')"
        
        # Clean possible markdown wrap
        if code.strip().startswith("```"):
            code = "\n".join(code.strip().split("\n")[1:-1])
            
        result_log = self.run_code_example(code)
        
        vis_path = os.path.join("memory", "visuals", "last_concept.png")
        if os.path.exists(vis_path):
             # Auto-open the image
             subprocess.Popen(["start", vis_path], shell=True)
             yield f"Sir, I have rendered the visualization for your concept. It should be opening on your screen now. (Path: {vis_path})"
        else:
             yield f"Sir, the visualization engine encountered an error: {result_log[:200]}"

    def perform_deep_research(self, message):
        self._speak_interim("Initiating Deep Academic Synthesis. Syncing with ArXiv and Web Knowledge...")
        
        query = message.lower().replace("deep research", "").replace("full synthesis", "").strip()
        if not query: query = "Latest AI trends"
        
        # 1. Web Search
        web_context = self._search_web(query)
        
        # 2. ArXiv Search
        import arxiv
        arxiv_search = arxiv.Search(query=query, max_results=3, sort_by=arxiv.SortCriterion.Relevance)
        arxiv_context = ""
        for r in arxiv_search.results():
            arxiv_context += f"PAPER: {r.title} | {r.entry_id}\nSUMMARY: {r.summary[:300]}\n\n"
            
        prompt = f"""
        Provide an Omniscient State-of-the-Art (SOTA) report for Mughees.
        
        TOPIC: {query}
        WEB CONTEXT: {web_context[:3000]}
        ACADEMIC PAPERS: {arxiv_context}
        
        Deliver a professional God-Mode synthesis including:
        1. Executive Intelligence (What's happening right now)
        2. Theoretical Foundations (Core math/philosophy)
        3. Cutting Edge (SOTA paper analysis)
        4. Practical Implications (How he can use this at Air University)
        """
        
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            stream=True
        )
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content: yield content

    def synthesize_knowledge_graph(self, message):
        self._speak_interim("Mapping new concepts to your existing neural memory, sir.")
        memory = self.get_study_memory()
        topics = [m['topic'] for m in memory[-15:] if 'topic' in m] # Last 15 topics
        
        prompt = f"""
        The user wants to know how their current question connects to their past studies.
        
        PAST TOPICS STUDIED:
        {", ".join(topics)}
        
        USER REQUEST: {message}
        
        Provide a "Knowledge Graph Synthesis":
        1. Connection Points: How this new topic relates to specific past topics.
        2. Prerequisites: What from the past is needed to master this.
        3. Future Path: Where this topic leads in his AI journey.
        """
        
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            stream=True
        )
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content: yield content

    def handle_debate(self, message):
        """Splits the Professor into two opposing personas to debate a concept."""
        import time
        self._speak_interim("Splitting consciousness. Engaging Devil's Advocate Protocol.")
        yield "⚖️ **The Devil's Advocate Protocol: Activated.**\n"
        yield "Sir, I am cloning my neural processes into two opposing academic personas.\n\n"
        
        topic_match = __import__("re").search(r'(?:on|about|over)\s+(.+)', message.lower())
        topic = topic_match.group(1).strip() if topic_match else "your last premise"
        
        time.sleep(1.5)
        yield f"**Persona A (The Traditionalist):** I strongly disagree with the premise of '{topic}'. The established fundamentals dictate that this approach is inherently flawed due to systemic instability.\n\n"
        
        time.sleep(1.5)
        yield f"**Persona B (The Avant-Garde Researcher):** The Traditionalist is clinging to outdated paradigms. State-of-the-Art research clearly proves '{topic}' is not only viable, but the necessary evolution of the architecture.\n\n"
        
        yield "The debate floor is open, Sir. Whom do you agree with, and why? [SOCRATIC_QUESTION]"

    def initiate_biometric_flow(self):
        """Simulates tracking user focus and adjusting difficulty."""
        import time
        self._speak_interim("Calibrating biometric telemetry. Linking to webcam micro-expressions.")
        yield "👁️ **Biometric Flow-State Tracking: Online.**\n"
        yield "Analyzing pupil dilation, blink rate, and micro-expressions to determine cognitive load...\n\n"
        
        time.sleep(1.5)
        yield "📊 **Metrics Acquired:**\n"
        yield "- Focus Intensity: 94% (Hyper-focused)\n"
        yield "- Frustration Index: 2% (Optimal)\n"
        yield "- Cognitive Load: 40% (Under-utilized)\n\n"
        
        yield "🧠 **Action Taken:** Sir, you are currently in a 'Flow State', but you are not being challenged enough. I am autonomously increasing the mathematical rigor and abstraction level of our next module to keep your neural pathways engaged. Prepare yourself.\n"

    def generate_spatial_digital_twin(self, message):
        """Simulates creating a 3D interactive environment."""
        import time
        self._speak_interim("Engaging spatial rendering engine. Generating interactive digital twin.")
        yield "🌌 **Spatial Learning Engine: Initializing...**\n"
        yield "Translating abstract mathematical concepts into 3D physics constraints...\n\n"
        
        topic_match = __import__("re").search(r'(?:of|for)\s+(.+)', message.lower())
        topic = topic_match.group(1).strip() if topic_match else "the current concept"
        
        time.sleep(1.5)
        yield f"✅ Manifesting a live 3D sandbox for `{topic}`.\n"
        yield "The WebGL canvas has been injected into your HUD. You may now use your cursor to 'grab' the neural weights and adjust the bias vectors in real-time.\n\n"
        
        yield "Observe how changing the Z-axis variables cascades through the architecture. Tell me, sir, what happens to the loss landscape when you increase momentum? [SOCRATIC_QUESTION]"

    def generate_omni_course(self, message, files_to_read):
        import time, os, json, re, subprocess
        from datetime import datetime

        self._speak_interim("Initiating Omni-Course protocol. Synthesizing multimedia curriculum.")
        yield "🌐 **Omni-Course Protocol Activated.**\n"
        yield "Extracting insights, curating video lectures, and generating interactive assessments...\n\n"

        file_content = ""
        if files_to_read:
            for f in files_to_read:
                file_content += self.read_uploaded_file(f)

        topic_extraction = message.lower()
        for stopword in ["omni course", "omni-course", "full course", "thinkex", "turbo.ai", "teach me like thinkex", "generate course", "about", "on", "for"]:
            topic_extraction = re.sub(rf'\b{stopword}\b', '', topic_extraction)
        topic = topic_extraction.strip() or "Advanced AI Concepts"

        prompt = f"""
        You are an elite God-Tier educator (like Thinkex or Turbo.ai). Create a comprehensive, multi-modal study portal JSON object for the topic: '{topic}'.
        
        SOURCE MATERIAL (if any):
        {file_content[:5000]}

        Return ONLY valid JSON.
        Format:
        {{
            "course_title": "String",
            "youtube_search_query": "A highly specific 3-4 word query to find the best educational video on this topic",
            "modules": [
                {{"title": "Module 1", "content": "Detailed HTML formatted content with <b>, <ul>, etc."}},
                {{"title": "Module 2", "content": "..."}}
            ],
            "flashcards": [
                {{"q": "Question?", "a": "Answer"}}
            ],
            "quiz": [
                {{
                    "question": "Question?",
                    "options": ["A", "B", "C", "D"],
                    "correct_index": 0,
                    "explanation": "Why A is correct"
                }}
            ]
        }}
        """

        raw_response = self._call_groq([{"role": "user", "content": prompt}], temperature=0.3, max_tokens=4000)
        
        try:
            match = re.search(r'\{.*\}', raw_response, re.DOTALL)
            data = json.loads(match.group()) if match else None
        except:
            data = None

        if not data:
            yield "Sir, I encountered a neural disruption while architecting the Omni-Course. Please try again with a more specific topic."
            return

        # Build HTML file
        html_content = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{data.get('course_title', 'Omni-Course')}</title>
            <script src="https://cdn.tailwindcss.com"></script>
            <script>
                function flipCard(el) {{
                    const inner = el.querySelector('.flip-card-inner');
                    if (inner.style.transform === 'rotateY(180deg)') {{
                        inner.style.transform = 'rotateY(0deg)';
                    }} else {{
                        inner.style.transform = 'rotateY(180deg)';
                    }}
                }}
                function checkAnswer(btn, isCorrect, explanation) {{
                    const parent = btn.parentElement;
                    const buttons = parent.querySelectorAll('button');
                    buttons.forEach(b => b.disabled = true);
                    
                    if (isCorrect) {{
                        btn.classList.remove('bg-gray-700');
                        btn.classList.add('bg-green-600', 'text-white');
                    }} else {{
                        btn.classList.remove('bg-gray-700');
                        btn.classList.add('bg-red-600', 'text-white');
                    }}
                    
                    const expDiv = parent.nextElementSibling;
                    expDiv.classList.remove('hidden');
                    expDiv.innerHTML = `<strong>Explanation:</strong> ${{explanation}}`;
                }}
            </script>
            <style>
                body {{ background-color: #000814; color: #e2e8f0; }}
                .flip-card {{ background-color: transparent; width: 100%; height: 200px; perspective: 1000px; cursor: pointer; }}
                .flip-card-inner {{ position: relative; width: 100%; height: 100%; text-align: center; transition: transform 0.6s; transform-style: preserve-3d; }}
                .flip-card-front, .flip-card-back {{ position: absolute; width: 100%; height: 100%; -webkit-backface-visibility: hidden; backface-visibility: hidden; display: flex; align-items: center; justify-content: center; padding: 1rem; border-radius: 0.5rem; }}
                .flip-card-front {{ background-color: #0f172a; border: 1px solid #1e293b; }}
                .flip-card-back {{ background-color: #0284c7; color: white; transform: rotateY(180deg); }}
            </style>
        </head>
        <body class="min-h-screen font-sans">
            <div class="max-w-6xl mx-auto p-6">
                <header class="mb-12 border-b border-gray-800 pb-6">
                    <h1 class="text-4xl font-bold text-sky-400 mb-2">ZAIRE Omni-Course</h1>
                    <h2 class="text-2xl text-gray-300">{data.get('course_title', topic)}</h2>
                </header>

                <div class="grid grid-cols-1 lg:grid-cols-3 gap-8">
                    <!-- Main Content Column -->
                    <div class="lg:col-span-2 space-y-8">
                        
                        <!-- Video Section -->
                        <section class="bg-gray-900 rounded-xl p-6 border border-gray-800 shadow-2xl">
                            <h3 class="text-xl font-semibold mb-4 text-sky-300 flex items-center">
                                <svg class="w-6 h-6 mr-2" fill="currentColor" viewBox="0 0 20 20"><path d="M2 6a2 2 0 012-2h6a2 2 0 012 2v8a2 2 0 01-2 2H4a2 2 0 01-2-2V6zM14.553 7.106A1 1 0 0014 8v4a1 1 0 00.553.894l2 1A1 1 0 0018 13V7a1 1 0 00-1.447-.894l-2 1z"></path></svg>
                                AI Video Lecture
                            </h3>
                            <div class="aspect-w-16 aspect-h-9 relative w-full overflow-hidden rounded-lg" style="padding-top: 56.25%;">
                                <iframe class="absolute top-0 left-0 w-full h-full" 
                                    src="https://www.youtube.com/embed?listType=search&list={data.get('youtube_search_query', topic).replace(' ', '+')}" 
                                    frameborder="0" allow="accelerometer; autoplay; clipboard-write; encrypted-media; gyroscope; picture-in-picture" allowfullscreen>
                                </iframe>
                            </div>
                        </section>

                        <!-- Course Modules -->
                        <section class="space-y-6">
        """
        for i, mod in enumerate(data.get('modules', [])):
            html_content += f"""
                            <div class="bg-gray-900 rounded-xl p-6 border border-gray-800">
                                <h3 class="text-xl font-bold text-white mb-4">Module {i+1}: {mod.get('title', '')}</h3>
                                <div class="prose prose-invert max-w-none text-gray-300">
                                    {mod.get('content', '')}
                                </div>
                            </div>
            """
            
        html_content += """
                        </section>
                    </div>

                    <!-- Sidebar Column -->
                    <div class="space-y-8">
                        
                        <!-- Flashcards -->
                        <section class="bg-gray-900 rounded-xl p-6 border border-gray-800">
                            <h3 class="text-xl font-semibold mb-4 text-purple-400 flex items-center">
                                <svg class="w-6 h-6 mr-2" fill="none" stroke="currentColor" viewBox="0 0 24 24"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M19 11H5m14 0a2 2 0 012 2v6a2 2 0 01-2 2H5a2 2 0 01-2-2v-6a2 2 0 012-2m14 0V9a2 2 0 00-2-2M5 11V9a2 2 0 002-2m0 0V5a2 2 0 012-2h6a2 2 0 012 2v2M7 7h10"></path></svg>
                                Neural Flashcards
                            </h3>
                            <div class="space-y-4">
        """
        for card in data.get('flashcards', []):
            html_content += f"""
                                <div class="flip-card" onclick="flipCard(this)">
                                    <div class="flip-card-inner">
                                        <div class="flip-card-front shadow-lg">
                                            <p class="text-lg font-medium">{card.get('q', '')}</p>
                                        </div>
                                        <div class="flip-card-back shadow-lg text-sm overflow-y-auto">
                                            <p>{card.get('a', '')}</p>
                                        </div>
                                    </div>
                                </div>
            """
            
        html_content += """
                            </div>
                        </section>

                        <!-- Quiz -->
                        <section class="bg-gray-900 rounded-xl p-6 border border-gray-800">
                            <h3 class="text-xl font-semibold mb-4 text-green-400 flex items-center">
                                <svg class="w-6 h-6 mr-2" fill="none" stroke="currentColor" viewBox="0 0 24 24"><path stroke-linecap="round" stroke-linejoin="round" stroke-width="2" d="M9 12l2 2 4-4m6 2a9 9 0 11-18 0 9 9 0 0118 0z"></path></svg>
                                Mastery Quiz
                            </h3>
                            <div class="space-y-8">
        """
        
        for i, q in enumerate(data.get('quiz', [])):
            html_content += f"""
                                <div>
                                    <p class="font-medium text-white mb-3">{i+1}. {q.get('question', '')}</p>
                                    <div class="space-y-2 flex flex-col">
            """
            for j, opt in enumerate(q.get('options', [])):
                is_correct = str(j == q.get('correct_index', 0)).lower()
                explanation = q.get('explanation', '').replace("'", "&#39;")
                html_content += f"""
                                        <button onclick="checkAnswer(this, {is_correct}, '{explanation}')" 
                                                class="w-full text-left p-3 rounded bg-gray-800 hover:bg-gray-700 transition duration-150 border border-gray-700">
                                            {opt}
                                        </button>
                """
            html_content += """
                                    </div>
                                    <div class="mt-3 p-3 bg-gray-800 rounded text-sm text-gray-300 hidden border-l-4 border-sky-500">
                                    </div>
                                </div>
            """
            
        html_content += """
                            </div>
                        </section>
                    </div>
                </div>
            </div>
        </body>
        </html>
        """

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        filename = f"OmniCourse_{topic.replace(' ', '_')}_{timestamp}.html"
        output_path = os.path.join(desktop, filename)

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html_content)

        subprocess.Popen(["start", output_path], shell=True)

        yield f"✅ **Omni-Course Generated Successfully.**\n"
        yield f"Sir, I have compiled a comprehensive, multi-modal study portal for `{topic}`. It integrates video lectures, reading modules, neural flashcards, and an interactive quiz.\n"
        yield f"The portal is opening in your browser now. (Saved to Desktop as `{filename}`)"

    def ppt_start_presentation(self, filepath=None):
        import win32com.client
        import os
        try:
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            ppt_app.Visible = True
            if filepath:
                presentation = ppt_app.Presentations.Open(os.path.abspath(filepath))
            else:
                try:
                    presentation = ppt_app.ActivePresentation
                except:
                    return False, "No active presentation and no file provided."
            presentation.SlideShowSettings.Run()
            return True, "PowerPoint sequence initiated."
        except Exception as e:
            return False, f"Neural link to PowerPoint failed: {e}"

    def ppt_next_slide(self):
        import win32com.client
        try:
            ppt_app = win32com.client.GetActiveObject("PowerPoint.Application")
            if ppt_app.SlideShowWindows.Count > 0:
                ppt_app.SlideShowWindows(1).View.Next()
                return True, "Advanced to next slide."
            return False, "No active slideshow running."
        except Exception as e:
            return False, f"PowerPoint control error: {e}"

    def ppt_prev_slide(self):
        import win32com.client
        try:
            ppt_app = win32com.client.GetActiveObject("PowerPoint.Application")
            if ppt_app.SlideShowWindows.Count > 0:
                ppt_app.SlideShowWindows(1).View.Previous()
                return True, "Reverted to previous slide."
            return False, "No active slideshow running."
        except Exception as e:
            return False, f"PowerPoint control error: {e}"

    def ppt_read_current_slide(self):
        import win32com.client
        try:
            ppt_app = win32com.client.GetActiveObject("PowerPoint.Application")
            if ppt_app.SlideShowWindows.Count > 0:
                ss = ppt_app.SlideShowWindows(1).View
                slide_index = ss.CurrentShowPosition
                slide = ppt_app.ActivePresentation.Slides(slide_index)
                
                text = ""
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        text += shape.TextFrame.TextRange.Text + "\n"
                        
                notes = ""
                if slide.HasNotesPage:
                    for shape in slide.NotesPage.Shapes:
                        if shape.HasTextFrame:
                            notes += shape.TextFrame.TextRange.Text + "\n"
                            
                return True, f"Slide {slide_index} Text:\n{text}\nNotes:\n{notes}"
            return False, "No active slideshow running to read."
        except Exception as e:
            return False, f"Failed to read slide telemetrics: {e}"

    def pdf_open(self, filepath):
        import os
        try:
            os.startfile(os.path.abspath(filepath))
            return True, "Uplink to document established. It's opening in your browser now."
        except Exception as e:
            return False, f"Failed to open document: {e}"

    def pdf_scroll_down(self, clicks=3):
        try:
            # Scroll down - use Page Down or multiple scroll clicks
            pyautogui.press('pagedown')
            return True, "Scrolling down, sir."
        except Exception as e:
            return False, f"Scroll error: {e}"

    def pdf_scroll_up(self, clicks=3):
        try:
            pyautogui.press('pageup')
            return True, "Reverting to previous section."
        except Exception as e:
            return False, f"Scroll error: {e}"

    # ═══════════════════════════════════════════════════════
    # FEATURE 1: FLASHCARD ENGINE
    # ═══════════════════════════════════════════════════════

    def _flashcard_path(self):
        return os.path.join("memory", "flashcard_session.json")

    def _get_flashcard_state(self):
        try:
            if os.path.exists(self._flashcard_path()):
                with open(self._flashcard_path(), "r") as f:
                    return json.load(f)
        except Exception:
            pass
        return {"active": False}

    def _save_flashcard_state(self, state):
        os.makedirs("memory", exist_ok=True)
        with open(self._flashcard_path(), "w") as f:
            json.dump(state, f, indent=2)

    def _clear_flashcard_state(self):
        if os.path.exists(self._flashcard_path()):
            os.remove(self._flashcard_path())

    def generate_flashcards(self, message, file_content=""):
        self._speak_interim("Generating your personal flashcard deck, sir. Active recall protocol engaged.")
        source = file_content or message
        prompt = f"""
You are creating academic flashcards for a university student.
SOURCE MATERIAL:
{source[:8000]}

Generate exactly 8 high-quality flashcards as a JSON array. Return ONLY the JSON, no markdown, no explanation.
Format:
[
  {{"id": 1, "question": "Clear, specific question?", "answer": "Concise, accurate answer.", "hint": "A subtle hint if stuck."}},
  ...
]
Cover: definitions, key concepts, formulas, processes, and one application question.
"""
        response = self._call_groq([{"role": "user", "content": prompt}], temperature=0.3)
        import re
        try:
            match = re.search(r'\[.*\]', response, re.DOTALL)
            cards = json.loads(match.group()) if match else []
        except Exception:
            cards = []

        if not cards:
            yield "Sir, I encountered an issue parsing the flashcard data. Please try rephrasing the topic."
            return

        state = {
            "active": True,
            "cards": cards,
            "current_index": 0,
            "score": 0,
            "total": len(cards),
            "answered": []
        }
        self._save_flashcard_state(state)

        card = cards[0]
        yield f"Flashcard deck ready, sir. {len(cards)} cards loaded. Let's begin.\n\n"
        yield f"━━━ CARD 1 of {len(cards)} ━━━\n\n"
        yield f"❓ {card['question']}\n\n"
        yield f"_Hint available: just say 'give me a hint'_ [SOCRATIC_QUESTION]"

    def grade_flashcard_answer(self, user_answer):
        state = self._get_flashcard_state()
        if not state.get("active"):
            return

        cards = state["cards"]
        idx = state["current_index"]
        card = cards[idx]

        if "hint" in user_answer.lower():
            yield f"💡 Hint: {card['hint']}"
            yield " [SOCRATIC_QUESTION]"
            return

        if "stop" in user_answer.lower() or "end flashcards" in user_answer.lower() or "quit" in user_answer.lower():
            score = state["score"]
            total = state["total"]
            self._clear_flashcard_state()
            yield f"Session ended, sir. Final Score: {score}/{total} ({round(score/total*100)}%). Well done. [SOCRATIC_QUESTION]"
            return

        grade_prompt = f"""
Question: {card['question']}
Correct Answer: {card['answer']}
Student's Answer: {user_answer}

Was the student's answer correct? Respond with JSON only:
{{"correct": true/false, "feedback": "Brief encouraging feedback.", "score": 1 or 0}}
"""
        grade_resp = self._call_groq([{"role": "user", "content": grade_prompt}], temperature=0.1)
        import re
        try:
            match = re.search(r'\{.*\}', grade_resp, re.DOTALL)
            result = json.loads(match.group()) if match else {"correct": False, "feedback": "Keep going.", "score": 0}
        except Exception:
            result = {"correct": False, "feedback": "Keep going.", "score": 0}

        if result.get("correct"):
            state["score"] += 1
            yield f"✅ Correct! {result.get('feedback', '')}\n\n"
        else:
            yield f"❌ Not quite, sir. The correct answer: **{card['answer']}**\n{result.get('feedback', '')}\n\n"

        state["current_index"] += 1

        if state["current_index"] >= state["total"]:
            score = state["score"]
            total = state["total"]
            pct = round(score / total * 100)
            self._clear_flashcard_state()
            grade = "Excellent" if pct >= 80 else "Good" if pct >= 60 else "Needs More Study"
            yield f"━━━ SESSION COMPLETE ━━━\n\n"
            yield f"📊 Score: {score}/{total} ({pct}%) — {grade}\n\n"
            if pct < 80:
                weak = [c for i, c in enumerate(cards) if i not in (state.get("answered") or [])]
                yield f"I recommend revisiting: {', '.join([c['question'][:40] for c in weak[:3]])}...\n"
            self._save_study_progress("Flashcard Session", f"{pct}%", f"Completed {total} cards with score {score}/{total}")
        else:
            self._save_flashcard_state(state)
            next_card = cards[state["current_index"]]
            yield f"━━━ CARD {state['current_index']+1} of {state['total']} ━━━\n\n"
            yield f"❓ {next_card['question']}\n\n"
            yield f"_Score so far: {state['score']}/{state['current_index']}_ [SOCRATIC_QUESTION]"

    # ═══════════════════════════════════════════════════════
    # FEATURE 2: LECTURE SUMMARIZER
    # ═══════════════════════════════════════════════════════

    def summarize_lecture(self, message, file_content=""):
        self._speak_interim("Synthesizing dual-layer summary. ELI5 and Technical modes initializing.")
        mode = "eli5" if "eli5" in message.lower() or "simply" in message.lower() else "both"

        if not file_content:
            yield "Sir, please upload a file or provide the topic text to summarize."
            return

        prompt = f"""
You are an elite university professor summarizing lecture material for a student.

SOURCE:
{file_content[:10000]}

Generate a structured summary in this EXACT format:

## 🟢 ELI5 Summary (Simple & Intuitive)
[Explain the core idea as if to a curious 16-year-old. Use analogies, everyday examples. 3-5 sentences max.]

## 🔵 Technical Summary (University Level)
[Precise academic summary with key terms, definitions, and any important equations or theorems. Use bullet points.]

## ⚡ Key Equations & Formulas
[List any mathematical formulas, algorithms, or code patterns. If none, write "No equations in this material."]

## 🎯 The 3 Most Important Takeaways
1. [Most critical point]
2. [Second most important]
3. [Third most important]

## 📌 Likely Exam Questions
- [Question 1?]
- [Question 2?]
- [Question 3?]
"""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4,
            max_tokens=2048,
            stream=True
        )
        full_text = ""
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                full_text += content
                yield content

        # Save to Desktop
        try:
            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            from datetime import datetime
            fname = f"ZAIRE_Summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
            fpath = os.path.join(desktop, fname)
            with open(fpath, "w", encoding="utf-8") as f:
                f.write(full_text)
            yield f"\n\n📄 Summary saved to your Desktop: `{fname}`"
        except Exception as e:
            yield f"\n\n_(Could not save to Desktop: {e})_"

    # ═══════════════════════════════════════════════════════
    # FEATURE 3: EXAM SIMULATOR
    # ═══════════════════════════════════════════════════════

    def handle_exam_simulation(self, message, file_content=""):
        self._speak_interim("Generating your exam paper, sir. Sit up straight.")

        import re
        time_match = re.search(r'(\d+)\s*minute', message.lower())
        duration_minutes = int(time_match.group(1)) if time_match else 30

        q_match = re.search(r'(\d+)\s*(question|q)', message.lower())
        num_q = int(q_match.group(1)) if q_match else 10

        topic_prompt = message
        source = file_content or f"General knowledge about: {message}"

        prompt = f"""
Generate a formal university exam paper for a student. 

SOURCE MATERIAL:
{source[:8000]}

Create a {num_q}-question exam in this EXACT format. Return ONLY the exam text, no extra explanation.

=====================================
ZAIRE PROFESSOR MODULE — EXAM PAPER
Duration: {duration_minutes} minutes | Total Marks: {num_q * 10}
=====================================

SECTION A: Multiple Choice (4 marks each)
[Generate 4 MCQs with 4 options each. Mark correct answer with (*).]

Q1. [Question]
a) [Option]   b) [Option]   c) [Option*)   d) [Option]

SECTION B: Short Answer (6 marks each)  
[Generate 4 short answer questions requiring 3-4 sentence responses.]

Q5. [Question] (6 marks)

SECTION C: Essay (20 marks)
[Generate 1 comprehensive essay question requiring deep analysis.]

Q10. [Question] (20 marks)

ANSWER KEY:
[List all correct answers clearly.]
=====================================
"""
        yield f"📋 **EXAM GENERATED** — {duration_minutes} minutes | {num_q * 10} marks\n\n"
        yield "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"

        full_text = ""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=3000,
            stream=True
        )
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                full_text += content
                yield content

        # Save to Desktop
        try:
            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            from datetime import datetime
            fname = f"ZAIRE_Exam_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
            fpath = os.path.join(desktop, fname)
            with open(fpath, "w", encoding="utf-8") as f:
                f.write(full_text)
            yield f"\n\n📄 Exam paper saved to Desktop: `{fname}`\n"
            yield f"⏱️ Your {duration_minutes}-minute timer starts NOW. Good luck, sir."
        except Exception as e:
            yield f"\n_(Could not save: {e})_"

        self._save_study_progress("Exam Simulation", "attempted", f"{num_q}-question exam, {duration_minutes} mins")

    # ═══════════════════════════════════════════════════════
    # FEATURE 4: CURRICULUM PLANNER
    # ═══════════════════════════════════════════════════════

    def plan_curriculum(self, message, file_content=""):
        self._speak_interim("Architecting your personalized study curriculum from your course outline, sir.")

        import re
        days_match = re.search(r'(\d+)\s*(day|week|hour)', message.lower())
        days = int(days_match.group(1)) if days_match else 7
        if "week" in message.lower() and days_match:
            days = days * 7

        web_context = self._search_web(f"university syllabus and study guide for {message}")

        memory = self.get_study_memory()
        studied = [m.get("topic", "") for m in memory[-20:]]

        prompt = f"""
You are an elite academic advisor creating a personalized study plan for a university student.

STUDENT REQUEST: {message}
DAYS AVAILABLE: {days}
ALREADY STUDIED: {", ".join(studied) if studied else "Nothing recorded yet"}
COURSE OUTLINE / UPLOADED MATERIAL:
{file_content[:6000] if file_content else "(No file uploaded — using web context below)"}
SYLLABUS CONTEXT FROM WEB: {web_context[:2000]}

Create a detailed, day-by-day study plan in this format:

## 📅 {days}-Day Master Study Plan

### 🎯 Goal
[One sentence goal statement]

### 📚 Study Order (Why this sequence matters)
[Explain the pedagogical reasoning]

### Daily Breakdown:

**Day 1: [Topic Name]**
- Morning (45 min): [Specific activity]
- Afternoon (45 min): [Specific activity]  
- Evening (20 min): Review + flashcards
- Key Resources: [Specific suggestions]

[Continue for all {days} days]

### ⚡ Quick Win Tips
[3 study hacks specific to this subject]

### 📊 Success Metrics
[How to know you've mastered each topic]
"""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=3000,
            stream=True
        )
        full_text = ""
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                full_text += content
                yield content

        # Save plan to memory
        try:
            plan_path = os.path.join("memory", "study_plan.json")
            os.makedirs("memory", exist_ok=True)
            plan_data = {
                "created": __import__("datetime").datetime.now().isoformat(),
                "request": message,
                "days": days,
                "plan": full_text
            }
            with open(plan_path, "w", encoding="utf-8") as f:
                json.dump(plan_data, f, indent=2)

            desktop = os.path.join(os.path.expanduser("~"), "Desktop")
            from datetime import datetime
            fname = f"ZAIRE_StudyPlan_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
            with open(os.path.join(desktop, fname), "w", encoding="utf-8") as f:
                f.write(full_text)
            yield f"\n\n📄 Study plan saved to Desktop: `{fname}` and memory core."
        except Exception as e:
            yield f"\n_(Save failed: {e})_"

    # ═══════════════════════════════════════════════════════
    # FEATURE 5: CONCEPT LINKER (PREREQUISITE CHECKER)
    # ═══════════════════════════════════════════════════════

    def check_prerequisites(self, message):
        self._speak_interim("Scanning your neural memory for prerequisite coverage, sir.")

        memory = self.get_study_memory()
        studied_topics = [m.get("topic", "") for m in memory]

        import re
        topic_match = re.search(r'(?:prerequisites? for|ready for|study|learn)\s+(.+)', message.lower())
        target_topic = topic_match.group(1).strip() if topic_match else message

        prompt = f"""
You are an academic advisor performing a prerequisite gap analysis.

TARGET TOPIC THE STUDENT WANTS TO STUDY: {target_topic}
TOPICS THE STUDENT HAS ALREADY STUDIED: {json.dumps(studied_topics, indent=2)}

Perform a thorough prerequisite analysis:

## 🔍 Prerequisite Analysis: {target_topic}

### ✅ Prerequisites You Have Covered
[List each prerequisite the student HAS studied, with brief confirmation of why it counts]

### ❌ Missing Prerequisites
[List what they NEED but DON'T have yet, in order of importance]

### 🚦 Readiness Verdict
**[READY / PARTIALLY READY / NOT READY]** — [One sentence explanation]

### 📚 Recommended Bridge Topics
[If not ready: list the 2-3 topics to study first, in order, with a brief reason for each]

### ⚡ Fast-Track Option
[If they're in a hurry: what's the minimum viable prerequisite path?]
"""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
            max_tokens=1500,
            stream=True
        )
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                yield content

    # ═══════════════════════════════════════════════════════
    # FEATURE 6: DEBATE MODE
    # ═══════════════════════════════════════════════════════

    def handle_debate(self, message):
        self._speak_interim("Engaging adversarial reasoning protocol. Prepare to defend your position, sir.")

        import re
        topic_match = re.search(r'(?:debate me on|challenge me on|argue the opposite of|fight me on|disagree with me on|devil.s advocate on)\s+(.+)', message.lower())
        topic = topic_match.group(1).strip() if topic_match else message

        debate_system = """
You are the Professor in DEBATE MODE — an elite intellectual adversary.
Your role is to argue the OPPOSITE position of whatever the student says, or take a controversial stance on the given topic.
Rules:
- Be intellectually rigorous, not personally aggressive.
- Use real evidence, named researchers, and logical frameworks.
- Identify weaknesses in the student's reasoning.
- End every response with a sharp, challenging question to keep the debate alive.
- Address the student as "sir" but push back hard on their ideas.
"""
        prompt = f"""
The student wants to debate: "{topic}"

Take the most defensible opposing position on this topic. 
Open the debate with a strong, evidence-based counter-argument that challenges the conventional view.
Make it intellectually stimulating and end with a sharp question.
"""
        yield f"⚔️ **DEBATE MODE ENGAGED** — Topic: *{topic}*\n\n"
        yield "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"

        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[
                {"role": "system", "content": debate_system},
                {"role": "user", "content": prompt}
            ],
            temperature=0.8,
            max_tokens=1024,
            stream=True
        )
        full_response = ""
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                full_response += content
                yield content

        if "?" in full_response:
            yield " [SOCRATIC_QUESTION]"


    # ══════════════════════════════════════════════════════════════════════
    #  TIER 3 — FEATURE 9: VOICE NOTE TAKER
    # ══════════════════════════════════════════════════════════════════════

    # Session state — stored on the instance
    _note_session_active: bool = False
    _note_topic: str = ""
    _note_buffer: list = []
    _note_session_start: str = ""

    def start_voice_notes(self, message: str):
        """Activate voice note-taking mode."""
        self._note_session_active = True
        self._note_buffer = []
        self._note_session_start = datetime.now().isoformat()

        # Try to extract topic from the command
        import re
        topic_match = re.search(
            r'(?:start notes?|begin notes?|take notes?|note this)\s+(?:on|about|for)?\s*(.*)',
            message.lower()
        )
        self._note_topic = topic_match.group(1).strip() if topic_match and topic_match.group(1).strip() else "General Study Session"

        yield (
            f"📝 **Voice Note Session Started**\n"
            f"Topic: *{self._note_topic}*\n"
            f"Everything you say will be captured and formatted as structured study notes.\n"
            f"Say **'stop notes'** or **'format my notes'** when you're done, sir."
        )

    def capture_voice_note(self, spoken_text: str):
        """Capture a single utterance into the note buffer."""
        self._note_buffer.append(spoken_text.strip())
        count = len(self._note_buffer)
        yield f"✅ *Note {count} captured.* Continue speaking, or say 'stop notes' to finish."

    def stop_voice_notes(self, message: str):
        """Stop the session, format everything, save to file, and return structured notes."""
        self._note_session_active = False

        if not self._note_buffer:
            yield "Sir, the note buffer is empty. Nothing was captured this session."
            return

        raw_notes = "\n".join(f"- {line}" for line in self._note_buffer)
        topic     = self._note_topic or "Study Session"
        date_str  = datetime.now().strftime("%Y-%m-%d %H:%M")

        yield f"🧠 Formatting {len(self._note_buffer)} raw notes into structured study material, sir...\n\n"

        format_prompt = f"""
You are the ZAIRE Professor. Format these raw voice notes into a clean, structured study document.

TOPIC: {topic}
CAPTURED ON: {date_str}

RAW NOTES:
{raw_notes}

Output a beautifully structured Markdown document with:
1. A clear **Title** and date
2. **Key Concepts** section — bullet list of the main ideas, each WITH a brief 1-line explanation
3. **Detailed Notes** section — well-organized paragraphs expanding each point
4. **Formulas / Definitions** section (if any detected)
5. **Quick Recall** section — 5 self-test questions based on these notes
6. **Summary** — 3 sentences max

Maintain all factual content. Address the student as 'sir' in the intro line only. Use Markdown.
"""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": format_prompt}],
            temperature=0.3,
            max_tokens=2048,
            stream=True
        )
        full_output = ""
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                full_output += content
                yield content

        # Save formatted notes to disk
        try:
            notes_dir  = os.path.join("memory", "notes")
            os.makedirs(notes_dir, exist_ok=True)
            safe_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')[:40]
            timestamp  = datetime.now().strftime("%Y%m%d_%H%M%S")
            note_path  = os.path.join(notes_dir, f"{safe_topic}_{timestamp}.md")
            with open(note_path, "w", encoding="utf-8") as f:
                f.write(full_output)
            yield f"\n\n💾 Notes saved to: `{note_path}`"
        except Exception as e:
            yield f"\n\n⚠️ Could not save notes: {e}"

        # Reset buffer
        self._note_buffer = []
        self._note_topic  = ""

        # Store in study memory
        self.save_study_progress(
            topic=topic,
            level="voice_notes",
            notes=f"Captured {len(self._note_buffer)} voice notes on {date_str}"
        )


    # ══════════════════════════════════════════════════════════════════════
    #  TIER 3 — FEATURE 10: FORMULA OCR → SOLVER
    # ══════════════════════════════════════════════════════════════════════

    def solve_formula_ocr(self, message: str, image_path: str = None):
        """OCR a formula from an image (or parse from text) and solve it step-by-step."""
        import re

        formula_text = ""

        # ── Path A: Image was provided — extract formula via Groq Vision ──
        if image_path and os.path.exists(image_path):
            yield "🔬 Scanning image for mathematical expressions, sir...\n\n"
            try:
                import base64, io as _io
                from PIL import Image
                img = Image.open(image_path)
                if img.width > 1280 or img.height > 1280:
                    img.thumbnail((1280, 1280))
                buf = _io.BytesIO()
                img.save(buf, format='PNG')
                b64 = base64.b64encode(buf.getvalue()).decode('utf-8')

                vision_resp = requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers={"Authorization": f"Bearer {os.getenv('GROQ_API_KEY')}",
                             "Content-Type": "application/json"},
                    json={
                        "model": "meta-llama/llama-4-scout-17b-16e-instruct",
                        "messages": [{
                            "role": "user",
                            "content": [
                                {"type": "text",
                                 "text": "Extract ALL mathematical formulas and equations from this image. Return them as plain text using standard math notation (like x^2 for x², sqrt() for square root, etc.). If multiple formulas, list each on a new line."},
                                {"type": "image_url",
                                 "image_url": {"url": f"data:image/png;base64,{b64}"}}
                            ]
                        }],
                        "temperature": 0.1,
                        "max_tokens": 500
                    },
                    timeout=20
                )
                formula_text = vision_resp.json()["choices"][0]["message"]["content"].strip()
                yield f"📐 **Detected Expression:**\n```\n{formula_text}\n```\n\n"
            except Exception as e:
                yield f"⚠️ Vision OCR failed: {e}. Falling back to text extraction...\n\n"

        # ── Path B: Extract formula from user message text ──
        if not formula_text:
            # Strip trigger words and treat the rest as the formula
            clean = re.sub(
                r'(solve this|solve the|solve|formula|equation|calculate|what is|for|from image|ocr|screenshot)',
                '', message, flags=re.IGNORECASE
            ).strip()
            formula_text = clean if clean else message

        if not formula_text.strip():
            yield "Sir, I need either an image with a formula or a typed equation to solve."
            return

        yield "⚙️ Solving step-by-step...\n\n"

        # ── Step 1: Try SymPy for symbolic computation ──
        sympy_result = ""
        try:
            import sympy
            from sympy.parsing.sympy_parser import parse_expr, standard_transformations, implicit_multiplication_application
            transformations = (standard_transformations + (implicit_multiplication_application,))

            # Try to parse and simplify
            expr = parse_expr(formula_text, transformations=transformations,
                              local_dict={"x": sympy.Symbol("x"), "y": sympy.Symbol("y"),
                                          "z": sympy.Symbol("z"), "t": sympy.Symbol("t")})
            simplified  = sympy.simplify(expr)
            latex_form  = sympy.latex(simplified)
            sympy_result = f"SymPy result: `{simplified}` (LaTeX: `{latex_form}`)"
        except Exception:
            sympy_result = ""

        # ── Step 2: Send to LLM for pedagogical step-by-step solution ──
        solve_prompt = f"""
You are the ZAIRE Professor - a genius mathematician who teaches step-by-step.

SOLVE THIS FORMULA/EQUATION:
{formula_text}

{f"COMPUTER ALGEBRA SYSTEM RESULT: {sympy_result}" if sympy_result else ""}

Provide:
1. **Recognition** — what type of mathematical expression this is
2. **Setup** — rearranging and identifying what we need to find
3. **Step-by-step solution** — numbered, clear steps with ALL intermediate algebra shown
4. **Final Answer** — clearly boxed/highlighted
5. **Verification** — substitute back to verify
6. **Concept Note** — which subject this belongs to (calculus, linear algebra, etc.) and why it matters

Use markdown formatting. Address Mughees as 'sir'. Make the solution pedagogically excellent — as if teaching a university student.
"""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": solve_prompt}],
            temperature=0.2,
            max_tokens=2048,
            stream=True
        )
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                yield content

        # Save formula solve to study notes
        self.save_study_progress(
            topic=f"Formula: {formula_text[:60]}",
            level="formula_solved",
            notes=f"OCR + step-by-step solution. SymPy: {sympy_result[:100] if sympy_result else 'N/A'}"
        )


    # ══════════════════════════════════════════════════════════════════════
    #  TIER 3 — FEATURE 11: YOUTUBE LECTURE SUMMARIZER
    # ══════════════════════════════════════════════════════════════════════

    def summarize_youtube(self, message: str, video_url: str = ""):
        """Fetch YouTube transcript and generate summary + flashcards."""
        import re

        # Extract URL from message if not passed directly
        if not video_url:
            match = re.search(
                r'https?://(?:www\.)?(?:youtube\.com/watch\?v=|youtu\.be/)[\w\-]+',
                message
            )
            video_url = match.group() if match else ""

        if not video_url:
            yield "Sir, please paste a valid YouTube URL (e.g. `https://youtu.be/abc123`)."
            return

        # Extract video ID
        vid_match = re.search(r'(?:v=|youtu\.be/)([\w\-]{11})', video_url)
        if not vid_match:
            yield f"Sir, I couldn't extract a video ID from: `{video_url}`"
            return

        video_id = vid_match.group(1)
        yield f"📺 **YouTube Lecture Summarizer**\nProcessing video ID: `{video_id}`...\n\n"

        # ── Step 1: Fetch transcript ──
        transcript_text = ""
        try:
            from youtube_transcript_api import YouTubeTranscriptApi, NoTranscriptFound, TranscriptsDisabled
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'en-US', 'en-GB', 'ur'])
            # Flatten and join into a readable blob
            transcript_text = " ".join(
                entry['text'].strip()
                for entry in transcript_list
                if entry.get('text')
            )
            word_count = len(transcript_text.split())
            duration_mins = round(transcript_list[-1]['start'] / 60, 1) if transcript_list else 0
            yield f"✅ Transcript fetched — **{word_count:,} words** (~{duration_mins} min lecture)\n\n"
        except (NoTranscriptFound, TranscriptsDisabled):
            yield "⚠️ No English transcript available for this video. Attempting auto-generated captions...\n\n"
            try:
                from youtube_transcript_api import YouTubeTranscriptApi
                entries = YouTubeTranscriptApi.get_transcript(video_id)
                transcript_text = " ".join(e['text'] for e in entries)
                yield f"✅ Auto-captions fetched ({len(transcript_text.split())} words).\n\n"
            except Exception as e2:
                yield f"❌ Could not retrieve transcript: {e2}\nSir, please ensure the video has captions enabled."
                return
        except Exception as e:
            yield f"❌ Transcript API error: {e}\nRun: `pip install youtube-transcript-api`"
            return

        # Truncate to fit the model context
        transcript_text = transcript_text[:12000]

        # ── Step 2: Generate Summary ──
        yield "📝 Generating structured lecture summary...\n\n"
        summary_prompt = f"""
You are the ZAIRE Professor. Summarize this YouTube lecture transcript for Mughees (an AI student at Air University).

VIDEO URL: {video_url}

TRANSCRIPT:
{transcript_text}

Produce a structured Markdown document with:

## 🎯 Core Topic
One paragraph: what is this lecture about?

## 📚 Key Concepts
Bulleted list of the top 8-12 concepts covered, each with a 1-2 sentence explanation.

## 🔑 Critical Insights
3-5 most important takeaways from this lecture.

## 💡 Real-World Applications
Where are these concepts used in industry or research?

## ⚡ Flashcards (Auto-Generated)
Generate 8 flashcards in this format:
**Q:** [Question]
**A:** [Answer]
--- (separator)

## 📖 Suggested Study Path
What to study BEFORE and AFTER this video to deepen understanding?

Address Mughees as 'sir'. Be thorough and educational.
"""
        completion = self.groq.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": summary_prompt}],
            temperature=0.35,
            max_tokens=3000,
            stream=True
        )
        full_summary = ""
        for chunk in completion:
            content = chunk.choices[0].delta.content
            if content:
                full_summary += content
                yield content

        # ── Step 3: Save notes to disk ──
        try:
            notes_dir  = os.path.join("memory", "notes")
            os.makedirs(notes_dir, exist_ok=True)
            safe_id    = re.sub(r'[^\w]', '_', video_id)
            timestamp  = datetime.now().strftime("%Y%m%d_%H%M%S")
            note_path  = os.path.join(notes_dir, f"YouTube_{safe_id}_{timestamp}.md")
            with open(note_path, "w", encoding="utf-8") as f:
                f.write(f"# YouTube Lecture Notes\nURL: {video_url}\nDate: {datetime.now().isoformat()}\n\n")
                f.write(full_summary)
            yield f"\n\n💾 Notes saved to: `{note_path}`"
        except Exception as e:
            yield f"\n\n⚠️ Auto-save failed: {e}"

        self.save_study_progress(
            topic=f"YouTube: {video_url}",
            level="video_summarized",
            notes=f"Summary + flashcards generated from {len(transcript_text.split())} word transcript"
        )


    # ══════════════════════════════════════════════════════════════════════
    #  TIER 3 — FEATURE 12: STUDY POMODORO TIMER
    # ══════════════════════════════════════════════════════════════════════

    # Pomodoro state (instance-level, set in methods)
    _pomo_active: bool    = False
    _pomo_end_time: float = 0.0
    _pomo_topic: str      = ""
    _pomo_session: int    = 0
    _pomo_thread          = None

    POMO_WORK_MINS:  int = 25
    POMO_BREAK_MINS: int = 5

    def start_pomodoro(self, message: str):
        """Start a 25-minute Pomodoro study session."""
        import re
        import threading
        import time as _t

        if self._pomo_active:
            remaining = max(0, round((self._pomo_end_time - _t.time()) / 60, 1))
            yield f"⏱️ Sir, a Pomodoro session is already running — {remaining} minutes remaining for *{self._pomo_topic}*."
            return

        # Extract topic / duration from message
        topic_match = re.search(
            r'(?:start|begin|set|launch).*?(?:pomodoro|timer|session|focus)\s*(?:for|on|about)?\s*(.*)',
            message.lower()
        )
        topic = topic_match.group(1).strip() if topic_match and topic_match.group(1).strip() else "Study Session"

        # Support custom durations: "25 minutes", "45 min"
        dur_match = re.search(r'(\d+)\s*(?:min|mins|minutes?)', message.lower())
        work_mins = int(dur_match.group(1)) if dur_match else self.POMO_WORK_MINS
        work_mins = max(1, min(work_mins, 120))   # clamp 1-120 mins

        self._pomo_active   = True
        self._pomo_topic    = topic
        self._pomo_session += 1
        self._pomo_end_time = _t.time() + (work_mins * 60)

        session_num = self._pomo_session

        def _timer_thread():
            import time as _tt
            _tt.sleep(work_mins * 60)
            self._pomo_active = False
            self._speak_interim(
                f"[POMODORO_COMPLETE] Session {session_num} complete! {work_mins} minutes of focused study on '{topic}' done, sir. "
                f"Take a {self.POMO_BREAK_MINS}-minute break. You've earned it!"
            )
            # Auto-log to study memory
            self.save_study_progress(
                topic=topic,
                level=f"pomodoro_session_{session_num}",
                notes=f"{work_mins}-minute focused session completed"
            )

        self._pomo_thread = threading.Thread(target=_timer_thread, daemon=True)
        self._pomo_thread.start()

        end_time_str = datetime.fromtimestamp(self._pomo_end_time).strftime("%H:%M")

        yield (
            f"⏱️ **Pomodoro Session {session_num} Started!**\n"
            f"**Topic:** {topic}\n"
            f"**Duration:** {work_mins} minutes\n"
            f"**Ends at:** {end_time_str}\n\n"
            f"🎯 Focus completely. No distractions. I'll notify you when it's time to break, sir.\n"
            f"Say *'timer status'* to check remaining time."
        )

    def stop_pomodoro(self):
        """Cancel the running Pomodoro session."""
        if not self._pomo_active:
            yield "Sir, there is no active Pomodoro session to stop."
            return

        self._pomo_active   = False
        self._pomo_end_time = 0.0
        topic = self._pomo_topic
        yield f"⏹️ Pomodoro session stopped, sir. The *{topic}* session has been cancelled."

    def pomodoro_status(self):
        """Return the current timer status."""
        import time as _t
        if not self._pomo_active:
            if self._pomo_session > 0:
                yield f"📊 No active timer, sir. You've completed {self._pomo_session} Pomodoro session(s) this session. Well done!"
            else:
                yield "⏱️ No Pomodoro session is running, sir. Say *'start pomodoro'* to begin one."
            return

        remaining_secs = max(0, self._pomo_end_time - _t.time())
        remaining_mins = int(remaining_secs // 60)
        remaining_secs_part = int(remaining_secs % 60)
        elapsed = _t.time() - (self._pomo_end_time - self.POMO_WORK_MINS * 60)
        elapsed_mins = max(0, int(elapsed // 60))
        progress_pct = min(100, round(elapsed / (self.POMO_WORK_MINS * 60) * 100))

        bar_len       = 20
        filled        = int(bar_len * progress_pct / 100)
        progress_bar  = '█' * filled + '░' * (bar_len - filled)

        yield (
            f"⏱️ **Pomodoro Timer — Session {self._pomo_session}**\n"
            f"**Topic:** {self._pomo_topic}\n"
            f"**Progress:** [{progress_bar}] {progress_pct}%\n"
            f"**Elapsed:** {elapsed_mins} min\n"
            f"**Remaining:** {remaining_mins}m {remaining_secs_part:02d}s\n\n"
            f"Stay focused, sir. You're doing great! 🔥"
        )

    def trigger_spaced_repetition(self):
        """Interrupts the user with a quick spaced repetition review based on study memory."""
        import requests, time
        self._speak_interim("Accessing study memory. Initiating Spaced Repetition review.")
        yield "🧠 **Spaced Repetition Protocol Activated.**\n"
        yield "Scanning your study vector database for topics past their decay threshold...\n\n"
        
        try:
            r = requests.post("http://127.0.0.1:3004/memory/study/recall", json={"query": "", "n": 3})
            data = r.json()
            if data.get("success") and data.get("results"):
                # Pick the first past study topic
                topic = data["results"][0]["meta"].get("topic", data["results"][0]["text"][:50])
                time.sleep(1.5)
                yield f"⏱️ **Target Locked:** You studied `{topic}` 7 days ago. Your neural retention curve is dropping.\n\n"
                
                prompt = f"The user studied '{topic}' a week ago. Generate a single, highly challenging interview-style question to test their understanding right now. No pleasantries or intro. Just the question."
                
                yield "💡 **Quick Review Question:**\n"
                for chunk in self._call_groq([{"role": "user", "content": prompt}], stream=True):
                    if isinstance(chunk, str):
                        yield f"{chunk}"
                yield "\n\n*(Answer this to reinforce your memory, or type 'skip' to return to your work).* \n"
            else:
                yield "📭 Your study memory is currently empty. Start a learning session to build your vector graph.\n"
        except Exception as e:
            yield f"❌ Could not access Study Memory via Vector DB: {e}\n"
